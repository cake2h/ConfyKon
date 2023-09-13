from pptx import Presentation
from pptx.util import Pt

conference_name = 'Студенческая научная конференция'


def edit_template(response: int, member_name: str, mentor_name: str, section: str, work_name: str) -> Presentation:
    path = 'templates/template_' + str(response) + '.pptx'
    template = Presentation(path)

    template.slides[0].shapes[2].text_frame.paragraphs[0].text = conference_name
    template.slides[0].shapes[2].text_frame.paragraphs[0].font.size = Pt(20)
    template.slides[0].shapes[2].text_frame.paragraphs[0].font.name = 'Montserrat'

    template.slides[0].shapes[1].text_frame.paragraphs[0].text = member_name
    template.slides[0].shapes[1].text_frame.paragraphs[0].font.size = Pt(11)
    template.slides[0].shapes[1].text_frame.paragraphs[0].font.name = 'Montserrat'

    template.slides[0].shapes[1].text_frame.paragraphs[2].text = 'за участие в секции '\
                                                                 + '"' + section + '"\nс докладом '\
                                                                 + '"' + work_name + '"\n\nНаучный руководитель: '\
                                                                 + mentor_name
    template.slides[0].shapes[1].text_frame.paragraphs[2].font.size = Pt(11)
    template.slides[0].shapes[1].text_frame.paragraphs[2].font.name = 'Montserrat'

    template.save("test.pptx")


edit_template(1, "Герасимов Константин", "Евгений Петрович", "Компьютерные науки", "Робот пылесос")
